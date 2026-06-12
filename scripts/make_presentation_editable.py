from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = Path(r"C:\Users\rbeyz\Desktop\AirPulse_Global_Project_Presentation_FULL_STORY.pptx")
OUTPUT = ROOT / "AirPulse_Global_Project_Presentation_FULL_STORY_EDITABLE.pptx"

BLUE = RGBColor(47, 99, 216)
PURPLE = RGBColor(98, 87, 214)
DARK = RGBColor(22, 28, 45)
MUTED = RGBColor(92, 99, 112)
LIGHT = RGBColor(244, 247, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(26, 140, 95)
AMBER = RGBColor(201, 123, 0)


def find_shape_by_text(slide, needle: str):
    for shape in slide.shapes:
        if hasattr(shape, "text") and needle in shape.text:
            return shape
    raise ValueError(f"Shape containing {needle!r} was not found.")


def set_shape_text(shape, text: str, *, size: int = 18, color: RGBColor = DARK, bold: bool = False, align=PP_ALIGN.LEFT):
    text_frame = shape.text_frame
    text_frame.clear()
    lines = text.split("\n")
    first = True
    for line in lines:
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        paragraph.alignment = align
        first = False


def add_placeholder(slide, *, left: float, top: float, width: float, height: float, title: str, note: str, accent: RGBColor):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = accent

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = note
    r2.font.size = Pt(12)
    r2.font.color.rgb = MUTED

    return box


def add_flow_box(slide, *, left: float, top: float, width: float, height: float, title: str, lines: list[str], accent: RGBColor):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()

    title_p = tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_r = title_p.add_run()
    title_r.text = title
    title_r.font.size = Pt(16)
    title_r.font.bold = True
    title_r.font.color.rgb = accent

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(10)
        r.font.color.rgb = MUTED

    return box


def style_banner(shape, fill: RGBColor, line: RGBColor | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()


def make_slide_5_architecture_room(prs: Presentation) -> None:
    slide = prs.slides[4]
    body1 = find_shape_by_text(slide, "Frontend/Backend:")
    body2 = find_shape_by_text(slide, "Core Engine:")
    body3 = find_shape_by_text(slide, "Cache Layer:")
    set_shape_text(body1, "Frontend / Backend\nSingle Streamlit app layer for dashboard flow, navigation and session context.", size=16)
    set_shape_text(body2, "Core Engine\nParallel API ingestion, normalization and forecast orchestration.", size=16)
    set_shape_text(body3, "Cache Layer\n`st.cache_data` protects performance and provider rate limits.", size=16)
    add_flow_box(
        slide,
        left=6.1,
        top=2.05,
        width=1.35,
        height=2.15,
        title="Inputs",
        lines=["WAQI live feed", "Open-Meteo fallback", "Tomorrow.io wind"],
        accent=BLUE,
    )
    add_flow_box(
        slide,
        left=7.58,
        top=2.05,
        width=1.35,
        height=2.15,
        title="Process",
        lines=["Normalize", "merge sources", "rank stations"],
        accent=PURPLE,
    )
    add_flow_box(
        slide,
        left=9.06,
        top=2.05,
        width=1.35,
        height=2.15,
        title="Intelligence",
        lines=["Forecast", "analytics", "health guidance"],
        accent=GREEN,
    )
    add_flow_box(
        slide,
        left=10.54,
        top=2.05,
        width=1.0,
        height=2.15,
        title="Output",
        lines=["Dashboard", "reports", "CSV / JPEG / PDF"],
        accent=AMBER,
    )

    for left in (7.33, 8.81, 10.29):
        connector = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(left),
            Inches(2.77),
            Inches(0.18),
            Inches(0.34),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = LIGHT
        connector.line.color.rgb = BLUE

    caption = slide.shapes.add_textbox(Inches(6.15), Inches(4.55), Inches(5.25), Inches(0.7))
    set_shape_text(
        caption,
        "Live source data is unified, enriched with wind context, then converted into forecast, analytics and export-ready user outputs.",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def make_slide_9_validation_room(prs: Presentation) -> None:
    slide = prs.slides[8]
    title = find_shape_by_text(slide, "Holt-Winters Strategy")
    intro = find_shape_by_text(slide, "When native provider forecasts")
    b1 = find_shape_by_text(slide, "• Seasonality:")
    b2 = find_shape_by_text(slide, "• Trend:")
    b3 = find_shape_by_text(slide, "• Confidence:")

    set_shape_text(title, "Fallback Forecast Logic", size=19, bold=True, color=BLUE)
    set_shape_text(intro, "Use this block for the method summary only.\nMove proof to a chart or metric table on the left.", size=14, color=MUTED)
    set_shape_text(b1, "Seasonality: captures repeated urban cycles.", size=14)
    set_shape_text(b2, "Trend: tracks directional pollution shifts.", size=14)
    set_shape_text(b3, "Confidence: show 90% band beside actual vs forecast.", size=14)

    add_placeholder(
        slide,
        left=0.7,
        top=2.0,
        width=4.6,
        height=3.9,
        title="Add Forecast Chart",
        note="Best use: actual vs forecast + confidence band + 1 small MAPE note",
        accent=PURPLE,
    )


def make_slide_10_screenshot_room(prs: Presentation) -> None:
    slide = prs.slides[9]
    desc1 = find_shape_by_text(slide, 'A high-level overview of city health.')
    desc2 = find_shape_by_text(slide, 'Rich geographic analysis.')
    set_shape_text(desc1, "Keep 1 dashboard screenshot here.\nOptional: add 2 short callouts only.", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    set_shape_text(desc2, "Keep 1 station map screenshot here.\nOptional: mark hotspot or station selection.", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_placeholder(
        slide,
        left=0.7,
        top=2.55,
        width=4.95,
        height=2.55,
        title="Dashboard Screenshot",
        note="AQI dial, wind context, health summary",
        accent=BLUE,
    )
    add_placeholder(
        slide,
        left=6.62,
        top=2.55,
        width=4.95,
        height=2.55,
        title="Map Screenshot",
        note="Station density, popup or selected hotspot",
        accent=GREEN,
    )


def make_slide_12_analytics_room(prs: Presentation) -> None:
    slide = prs.slides[11]
    statements = [
        find_shape_by_text(slide, "Meteorological Correlation:"),
        find_shape_by_text(slide, "Source Identification:"),
        find_shape_by_text(slide, "Forecast Validation:"),
        find_shape_by_text(slide, "Anomaly Detection:"),
    ]
    concise = [
        "Meteorological correlation\nWind speed vs pollutant stagnation",
        "Source identification\nTraffic-linked NO2 vs industrial SO2 signals",
        "Forecast validation\nActual vs estimate comparison on recent windows",
        "Anomaly detection\nUnexpected spikes beyond historical seasonality",
    ]
    for shape, text in zip(statements, concise):
        set_shape_text(shape, text, size=14)

    add_placeholder(
        slide,
        left=6.9,
        top=2.05,
        width=4.2,
        height=3.7,
        title="Add Validation Evidence",
        note="Chart, anomaly table, or forecast accuracy snapshot",
        accent=PURPLE,
    )


def make_slide_14_metrics_room(prs: Presentation) -> None:
    slide = prs.slides[13]
    outcome_box = find_shape_by_text(slide, "Outcome:")
    outcome_text = find_shape_by_text(slide, "Outcome: AirPulse")

    # Turn the bottom banner into an empty metrics area.
    outcome_box.left = Inches(0.75)
    outcome_box.top = Inches(5.05)
    outcome_box.width = Inches(10.5)
    outcome_box.height = Inches(1.15)
    style_banner(outcome_box, LIGHT, BLUE)
    set_shape_text(
        outcome_text,
        "Add quantified proof here\nExamples: MAPE, validated points, sample cities, screenshot of forecast comparison",
        size=15,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    outcome_text.left = Inches(1.0)
    outcome_text.top = Inches(5.35)
    outcome_text.width = Inches(10.0)
    outcome_text.height = Inches(0.6)


def make_slide_15_roadmap_room(prs: Presentation) -> None:
    slide = prs.slides[14]
    next_step_title = find_shape_by_text(slide, "Next Step")
    next_step_body = find_shape_by_text(slide, "Scale AirPulse")
    next_step_note = find_shape_by_text(slide, "Environmental intelligence designed")
    contact_line = find_shape_by_text(slide, "contact@airpulse.io")

    set_shape_text(next_step_title, "Roadmap + Limits", size=20, bold=True, color=BLUE)
    set_shape_text(
        next_step_body,
        "Roadmap\n- Alerts / notifications\n- Stronger validation benchmark\n- Institutional reporting view",
        size=14,
    )
    set_shape_text(
        next_step_note,
        "Limits\n- API dependency\n- optional wind key\n- coverage varies by city",
        size=13,
        color=AMBER,
    )
    set_shape_text(
        contact_line,
        "Use this left area for final demo link, GitHub, or presenter contact details.",
        size=12,
        color=MUTED,
    )


def main() -> None:
    prs = Presentation(str(SOURCE))
    make_slide_5_architecture_room(prs)
    make_slide_9_validation_room(prs)
    make_slide_10_screenshot_room(prs)
    make_slide_12_analytics_room(prs)
    make_slide_14_metrics_room(prs)
    make_slide_15_roadmap_room(prs)
    prs.save(str(OUTPUT))
    print(f"Saved editable copy to: {OUTPUT}")


if __name__ == "__main__":
    main()
