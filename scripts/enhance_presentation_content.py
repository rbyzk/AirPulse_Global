from __future__ import annotations

import csv
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "AirPulse_Global_Project_Presentation_FULL_STORY_EDITABLE.pptx"
OUTPUT = ROOT / "AirPulse_Global_Project_Presentation_FULL_STORY_REVIEWED.pptx"

BLUE = RGBColor(47, 99, 216)
DARK = RGBColor(22, 28, 45)
MUTED = RGBColor(92, 99, 112)
LIGHT = RGBColor(244, 247, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(26, 140, 95)
AMBER = RGBColor(201, 123, 0)


def find_shape_by_text(slide, needle: str):
    for shape in slide.shapes:
        if hasattr(shape, "text") and needle in shape.text:
            return shape
    raise ValueError(f"Shape containing {needle!r} was not found.")


def set_shape_text(shape, text: str, *, size: int = 18, color: RGBColor = DARK, bold: bool = False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        p.alignment = align
        first = False


def add_card(slide, *, left: float, top: float, width: float, height: float, title: str, lines: list[str], accent: RGBColor):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(2)

    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = accent

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    return card


def read_metrics() -> dict[str, str]:
    metrics_path = ROOT / "artifacts" / "offline_forecast_metrics_summary.csv"
    with metrics_path.open(encoding="utf-8-sig", newline="") as handle:
        rows = list(csv.DictReader(handle))

    pm25 = next(row for row in rows if row["pollutant"] == "pm25")

    lso_path = ROOT / "data" / "processed" / "station_expansion" / "leave_station_out_overall_leaderboard.csv"
    with lso_path.open(encoding="utf-8-sig", newline="") as handle:
        lso_rows = list(csv.DictReader(handle))
    best_lso = next(row for row in lso_rows if row["model_name"] == "gradient_boosting")

    validation_path = ROOT / "data" / "processed" / "forecast_validation.json"
    records = json.loads(validation_path.read_text(encoding="utf-8"))["records"]
    usable = []
    for rec in records:
        try:
            actual = float(rec["actual_value"])
            pred = float(rec["predicted_value"])
        except (KeyError, TypeError, ValueError):
            continue
        if math.isnan(actual) or math.isnan(pred):
            continue
        usable.append(abs(actual - pred))

    obs_count = len(list((ROOT / "data" / "processed" / "observation_cache").glob("*.csv")))

    return {
        "pm25_r2": f"{float(pm25['r2']):.2f}",
        "pm25_rmse": f"{float(pm25['rmse']):.2f}",
        "pm10_mape": f"{float(next(r for r in rows if r['pollutant'] == 'pm10')['mape']):.1f}%",
        "lso_rows": f"{int(float(best_lso['rows'])):,}",
        "lso_smape": f"{float(best_lso['smape']):.1f}%",
        "validation_records": str(len(records)),
        "usable_records": str(len(usable)),
        "validation_mae": f"{(sum(usable) / len(usable)):.2f}" if usable else "n/a",
        "obs_count": str(obs_count),
    }


def update_slide_9(prs: Presentation, metrics: dict[str, str]) -> None:
    slide = prs.slides[8]
    title = find_shape_by_text(slide, "Fallback Forecast Logic")
    intro = find_shape_by_text(slide, "Use this block for the method summary only.")
    seasonality = find_shape_by_text(slide, "Seasonality:")
    trend = find_shape_by_text(slide, "Trend:")
    confidence = find_shape_by_text(slide, "Confidence:")

    set_shape_text(title, "Forecasting Strategy & Validation", size=18, color=BLUE, bold=True)
    set_shape_text(
        intro,
        "Official WAQI daily forecast is used first.\nIf upstream coverage is missing, AirPulse falls back to a conservative time-series estimate.",
        size=13,
        color=DARK,
    )
    set_shape_text(
        seasonality,
        f"Offline PM2.5 champion: R² {metrics['pm25_r2']} and RMSE {metrics['pm25_rmse']}.",
        size=13,
    )
    set_shape_text(
        trend,
        f"Cross-station benchmark: {metrics['lso_rows']} rows with best sMAPE {metrics['lso_smape']}.",
        size=13,
    )
    set_shape_text(
        confidence,
        f"Live validation store: {metrics['validation_records']} records collected; {metrics['usable_records']} already comparable to actuals.",
        size=13,
    )

    add_card(
        slide,
        left=6.65,
        top=5.05,
        width=4.55,
        height=0.7,
        title="Why this matters",
        lines=["The deck now shows both model logic and real evaluation evidence instead of a placeholder."],
        accent=GREEN,
    )


def update_slide_14(prs: Presentation, metrics: dict[str, str]) -> None:
    slide = prs.slides[13]
    proof = find_shape_by_text(slide, "Add quantified proof here")
    set_shape_text(
        proof,
        "Evidence Snapshot\n"
        f"- 48 cached city / station histories support repeatable demos\n"
        f"- PM2.5 offline champion: R² {metrics['pm25_r2']} | RMSE {metrics['pm25_rmse']}\n"
        f"- Leave-station-out benchmark: {metrics['lso_rows']} rows | best sMAPE {metrics['lso_smape']}\n"
        f"- Validation store: {metrics['validation_records']} records captured for forecast follow-up",
        size=13,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )


def update_slide_15(prs: Presentation) -> None:
    slide = prs.slides[14]
    placeholder = find_shape_by_text(slide, "Use this left area for final demo link, GitHub, or presenter contact details.")
    set_shape_text(
        placeholder,
        "Demo focus\n"
        "- Live web app already available\n"
        "- Best story: city search -> forecast -> take action -> report export\n"
        "- Mention that wind enrichment is stronger when Tomorrow.io is configured",
        size=13,
        color=DARK,
    )

    questions = find_shape_by_text(slide, "Questions?")
    set_shape_text(questions, "Demo / Q&A", size=18, color=BLUE, bold=True)


def main() -> None:
    metrics = read_metrics()
    prs = Presentation(str(SOURCE))
    update_slide_9(prs, metrics)
    update_slide_14(prs, metrics)
    update_slide_15(prs)
    prs.save(str(OUTPUT))
    print(f"Saved reviewed deck to: {OUTPUT}")


if __name__ == "__main__":
    main()
