"""Generate polished Word and PowerPoint deliverables for AirPulse Global."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
import csv

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PPTInches, Pt as PPTPt


ROOT = Path(__file__).resolve().parent
DOCX_PATH = ROOT / "AirPulse_Global_Proje_Sunumu.docx"
PPTX_PATH = ROOT / "AirPulse_Global_Sunumu.pptx"

BLUE = RGBColor(47, 99, 216)
PURPLE = RGBColor(98, 87, 214)
DARK = RGBColor(22, 28, 45)
MUTED = RGBColor(92, 99, 112)
GREEN = RGBColor(26, 140, 95)
PPT_BLUE = PPTColor(47, 99, 216)
PPT_PURPLE = PPTColor(98, 87, 214)
PPT_DARK = PPTColor(22, 28, 45)
PPT_MUTED = PPTColor(92, 99, 112)
PPT_GREEN = PPTColor(26, 140, 95)
PPT_LIGHT = PPTColor(244, 247, 252)


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        return []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.DictReader(handle))


def project_stats() -> dict[str, int | str]:
    raw_city_files = list((ROOT / "data" / "raw").glob("*.csv"))
    scripts = list((ROOT / "scripts").glob("*.py"))
    notebooks = list((ROOT / "notebooks").glob("*.ipynb"))
    modules = list((ROOT / "src" / "airpulse").rglob("*.py"))
    stations = read_csv_rows(ROOT / "stations.csv")
    accepted = read_csv_rows(
        ROOT / "data" / "processed" / "station_expansion" / "station_candidates_accepted.csv"
    )
    metrics = read_csv_rows(ROOT / "artifacts" / "offline_forecast_metrics_summary.csv")
    return {
        "generated_at": datetime.now().strftime("%d.%m.%Y"),
        "raw_city_files": len(raw_city_files),
        "scripts": len(scripts),
        "notebooks": len(notebooks),
        "modules": len(modules),
        "stations": len(stations),
        "accepted_candidates": len(accepted),
        "forecast_metric_rows": len(metrics),
    }


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_bullet(doc: Document, text: str, *, level: int = 0) -> None:
    paragraph = doc.add_paragraph(style="List Bullet")
    if level:
        paragraph.paragraph_format.left_indent = Cm(0.75 * level)
    run = paragraph.add_run(text)
    run.font.size = Pt(11)
    run.font.color.rgb = DARK


def add_number(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="List Number")
    run = paragraph.add_run(text)
    run.font.size = Pt(11)
    run.font.color.rgb = DARK


def add_heading(doc: Document, text: str, level: int = 1) -> None:
    heading = doc.add_heading(text, level=level)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = heading.runs[0]
    run.font.color.rgb = BLUE if level == 1 else PURPLE
    run.font.bold = True


def add_paragraph(doc: Document, text: str, *, size: int = 11, color: RGBColor = DARK) -> None:
    paragraph = doc.add_paragraph()
    paragraph.paragraph_format.space_after = Pt(8)
    run = paragraph.add_run(text)
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color


def create_word_document(stats: dict[str, int | str]) -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2)
    section.right_margin = Cm(2)

    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Aptos"
    normal_style.font.size = Pt(11)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("AirPulse Global")
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLUE

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("Proje Tanıtım Dokümanı")
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.color.rgb = PURPLE

    summary = doc.add_paragraph()
    summary.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = summary.add_run(
        "Açık hava kalitesi verisini gerçek zamanlı, anlaşılır ve eyleme dönüştürülebilir içgörüye çeviren çevresel zeka platformu"
    )
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED

    chips = doc.add_paragraph()
    chips.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = chips.add_run("Canlı veri  |  Global istasyon ağı  |  Tahmin  |  Kişiselleştirilmiş aksiyon  |  Raporlama")
    run.font.size = Pt(10)
    run.font.color.rgb = GREEN
    run.bold = True

    date_p = doc.add_paragraph()
    date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = date_p.add_run(f"Hazırlanma tarihi: {stats['generated_at']}")
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED

    doc.add_page_break()

    add_heading(doc, "1. Yönetici Özeti")
    add_paragraph(
        doc,
        "AirPulse Global, dünya genelindeki hava kalitesi verilerini tek bir web uygulamasında toplayan, bu verileri rüzgar bağlamı ve tahmin katmanlarıyla zenginleştiren, ardından kullanıcıya anlaşılır öneriler ve paylaşılabilir çıktılar sunan tam yığın bir çevresel zeka çözümüdür.",
    )
    add_paragraph(
        doc,
        "Platformun temel amacı, ham çevre verisini sadece göstermek değil; kullanıcıların günlük kararlarını destekleyecek şekilde yorumlamak, sadeleştirmek ve aksiyona dönüştürmektir. Bu sayede AirPulse Global, hem teknik ekiplerin hem de son kullanıcıların rahatça anlayabileceği bir karar destek deneyimi sunar.",
    )

    info_table = doc.add_table(rows=1, cols=2)
    info_table.style = "Table Grid"
    headers = info_table.rows[0].cells
    headers[0].text = "Gösterge"
    headers[1].text = "Değer"
    set_cell_shading(headers[0], "DCE9FF")
    set_cell_shading(headers[1], "DCE9FF")
    for key, value in [
        ("Python modülü", stats["modules"]),
        ("Yardımcı script", stats["scripts"]),
        ("Araştırma notebook'u", stats["notebooks"]),
        ("Ham şehir veri dosyası", stats["raw_city_files"]),
        ("İzleme istasyonu kaydı", stats["stations"]),
        ("Kabul edilen yeni istasyon adayı", stats["accepted_candidates"]),
        ("Tahmin metrik özeti", stats["forecast_metric_rows"]),
    ]:
        row = info_table.add_row().cells
        row[0].text = str(key)
        row[1].text = str(value)

    add_heading(doc, "2. AirPulse Global Nedir?")
    add_paragraph(
        doc,
        "AirPulse Global, açık hava kalitesi kaynaklarından alınan gerçek zamanlı ölçümleri; harita, tahmin, öneri ve raporlama katmanlarıyla birleştiren modüler bir uygulamadır. Uygulama kullanıcıya yalnızca 'hava bugün kötü mü?' sorusunun cevabını vermez; aynı zamanda 'neden böyle?', 'yarın ne olabilir?' ve 'ben bugün ne yapmalıyım?' sorularını da yanıtlar.",
    )
    add_bullet(doc, "Gerçek zamanlı AQI, PM2.5, PM10, NO2, O3 ve istasyon düzeyi ölçümler sunar.")
    add_bullet(doc, "Rüzgar hızı ve yönüyle hava kalitesini birlikte yorumlayarak bağlam üretir.")
    add_bullet(doc, "Tahmin sayfasında resmi WAQI tahminlerini, yoksa güvenli geri dönüş modeliyle devam ettirir.")
    add_bullet(doc, "Take Action alanında sağlık ve sürdürülebilirlik odaklı günlük öneriler üretir.")
    add_bullet(doc, "PDF, JPEG sosyal kart ve CSV dışa aktarımlarıyla sonuçları paylaşılabilir hale getirir.")

    add_heading(doc, "3. Kullanıcıya Sağladığı Temel Değer")
    add_number(doc, "Veriyi sadeleştirir: Farklı kaynaklardan gelen teknik hava kalitesi verisini anlaşılır bir kullanıcı deneyimine dönüştürür.")
    add_number(doc, "Bağlam ekler: Rüzgar, konum ve istasyon bilgisiyle kirleticilerin yalnızca seviyesini değil etkisini de gösterir.")
    add_number(doc, "Aksiyon üretir: Günlük davranış önerileri, egzersiz planlama ve karbon ayak izi takibi sunar.")
    add_number(doc, "Güven verir: Resmi tahmin verisi ile geri dönüş modelini açık biçimde ayırır; kullanıcı ne gördüğünü bilir.")
    add_number(doc, "Paylaşımı kolaylaştırır: Tek tıkla profesyonel rapor ve sosyal kart çıktıları üretir.")

    add_heading(doc, "4. Platform Özellikleri")
    add_heading(doc, "4.1 Dashboard", level=2)
    add_paragraph(
        doc,
        "Dashboard, AirPulse deneyiminin başlangıç noktasıdır. Kullanıcı aradığı şehri seçer, anlık AQI ve kirletici kartlarını görür, rüzgar bağlamını inceler ve bölgesel filtrelerle şehirleri karşılaştırır. Open-Meteo sayesinde 30'dan fazla şehir paralel şekilde yüklenebilir.",
    )
    add_heading(doc, "4.2 Stations Map", level=2)
    add_paragraph(
        doc,
        "İstasyon haritası, şehir bazlı özetin ötesine geçerek ölçüm noktalarının dağılımını görünür kılar. Bounds tabanlı WAQI araması ve global tile katmanı sayesinde kullanıcılar farklı istasyonları tek tek inceleyebilir, hangi istasyonun Forecast sayfasına kaynak olacağını seçebilir.",
    )
    add_heading(doc, "4.3 Forecast", level=2)
    add_paragraph(
        doc,
        "Tahmin katmanı, önce resmi WAQI günlük tahminlerini kullanır. Bu veri mevcut değilse observed-history tabanlı geri dönüş yaklaşımı devreye girer. Böylece uygulama hem kaynak şeffaflığını korur hem de tahmin sürekliliğini kaybetmez.",
    )
    add_heading(doc, "4.4 Take Action", level=2)
    add_paragraph(
        doc,
        "Take Action alanı kişisel karar katmanıdır. Karbon ayak izi hesaplama, işe gidiş senaryolarını karşılaştırma, günlük davranış listesi ve streak takibi gibi araçlar sayesinde kullanıcı, çevresel veriyi doğrudan yaşam tarzı kararlarına bağlayabilir.",
    )
    add_heading(doc, "4.5 Reports ve Analytics", level=2)
    add_paragraph(
        doc,
        "Reports, uygulamadaki kişisel ve çevresel özetleri A4 PDF, JPEG ve CSV çıktılara dönüştürür. Analytics ise zaman desenleri, rüzgar dispersiyonu, karşılaştırmalı analiz ve anomali gözlemleri gibi daha analitik okumalara alan açar.",
    )

    add_heading(doc, "5. Nasıl Çalışır?")
    add_number(doc, "Kullanıcı şehir veya istasyon seçer.")
    add_number(doc, "WAQI ve Open-Meteo kaynaklarından canlı hava kalitesi verisi alınır.")
    add_number(doc, "Tomorrow.io varsa rüzgar hızı, yönü ve gust verisi eklenir.")
    add_number(doc, "Veriler cache katmanından geçerek performans ve kota uyumluluğu korunur.")
    add_number(doc, "Tahmin sayfası resmi WAQI forecast verisini kontrol eder.")
    add_number(doc, "Eksik durumda geri dönüş modeli devreye girer ve güven bandı ile sunum yapılır.")
    add_number(doc, "Take Action modülü sağlık ve sürdürülebilirlik önerilerini oluşturur.")
    add_number(doc, "Raporlama servisi sonuçları PDF, JPEG ve CSV formatında dışa aktarır.")

    add_heading(doc, "6. Teknoloji Yığını")
    tech_table = doc.add_table(rows=1, cols=3)
    tech_table.style = "Table Grid"
    head = tech_table.rows[0].cells
    head[0].text = "Bileşen"
    head[1].text = "Teknoloji"
    head[2].text = "Projede Rolü"
    for cell in head:
        set_cell_shading(cell, "DCE9FF")
    tech_rows = [
        ("Uygulama Katmanı", "Streamlit", "Çok sayfalı arayüz, session state, hızlı ürün deneyimi"),
        ("Veri Kaynağı", "WAQI / AQICN API", "Canlı hava kalitesi verileri, istasyon keşfi, resmi forecast beslemesi"),
        ("Global Fallback", "Open-Meteo", "Anahtarsız dashboard özeti ve çok şehirli paralel veri çekimi"),
        ("Rüzgar Katmanı", "Tomorrow.io", "Rüzgar hızı, yönü ve gust bağlamı"),
        ("Görselleştirme", "Plotly", "Etkileşimli grafikler ve tahmin görselleri"),
        ("Harita", "Folium + Leaflet", "İstasyon yoğunluğu ve etkileşimli popup'lar"),
        ("Raporlama", "ReportLab", "A4 PDF üretimi"),
        ("Sosyal Kart", "Pillow", "1080x566 görsel çıktı üretimi"),
        ("Dil / Mimari", "Python 3.11+", "Modüler servisler, veri işleme, tahmin ve entegrasyon katmanları"),
    ]
    for row_data in tech_rows:
        row = tech_table.add_row().cells
        for index, value in enumerate(row_data):
            row[index].text = value

    add_heading(doc, "7. Mimari ve Ürün Tasarımı")
    add_paragraph(
        doc,
        "AirPulse Global üç katmanlı düşünülmüştür: kullanıcıların etkileşime geçtiği runtime uygulama katmanı, araştırma ve modelleme katmanı, bir de raporlama, çeviri, yapılandırma ve gizli anahtar yönetimi gibi destek servisleri. Bu yapı hem geliştirme hızını hem de bakım kolaylığını artırır.",
    )
    add_bullet(doc, "Modüler Python paket yapısı sayesinde veri çekme, tahmin, raporlama ve UI sorumlulukları ayrılmıştır.")
    add_bullet(doc, "st.cache_data kullanımı; performans, kullanıcı deneyimi ve API rate-limit uyumu için kritik rol oynar.")
    add_bullet(doc, "ThreadPoolExecutor ile çok şehirli veri çekimi hızlandırılır.")
    add_bullet(doc, "Session state; şehir, istasyon, checklist ve kişisel seçimlerin sayfalar arasında korunmasını sağlar.")
    add_bullet(doc, "Uygulama, üçüncü taraf servis eksik olduğunda tamamen durmak yerine kontrollü fallback davranışı sergiler.")

    add_heading(doc, "8. Tahmin Yaklaşımı")
    add_paragraph(
        doc,
        "Tahmin katmanının en önemli tasarım ilkesi şeffaflıktır. Kullanıcı resmi sağlayıcı tahmini ile uygulamanın geri dönüş tahminini karıştırmamalıdır. Bu nedenle AirPulse önce WAQI forecast verisini kullanır, veri yoksa observed-history tabanlı güvenli modele geçer ve bunu arayüzde açıkça belirtir.",
    )
    add_bullet(doc, "Birincil yol: WAQI resmi günlük PM2.5, PM10 ve O3 forecast verisi")
    add_bullet(doc, "Geri dönüş yolu: Holt-Winters tabanlı zaman serisi tahmini")
    add_bullet(doc, "Sunum yaklaşımı: artan güven bandı ve model açıklaması")
    add_bullet(doc, "Fayda: veri sürekliliği korunurken güven kaybı yaşanmaz")

    add_heading(doc, "9. Raporlama ve Çıktılar")
    add_paragraph(
        doc,
        "Platform yalnızca ekranda bilgi göstermekle kalmaz; sonucu farklı paydaşlarla paylaşılabilir hale de getirir. Bu nedenle raporlama katmanı teknik doğruluk ile görsel netlik arasında denge kurar.",
    )
    add_bullet(doc, "A4 PDF: AQI banner, kirletici tablosu, rüzgar özeti, karbon ayak izi ve aksiyon skoru")
    add_bullet(doc, "JPEG sosyal kart: hızlı paylaşım için özet görsel")
    add_bullet(doc, "CSV export: hesaplanan değerlerin analiz veya raporlama için dışa aktarımı")

    add_heading(doc, "10. API Yapılandırması ve Operasyonel Notlar")
    add_paragraph(
        doc,
        "Uygulama, WAQI anahtarı ile tam canlı deneyime geçer. Tomorrow.io anahtarı ise isteğe bağlıdır; yokluğunda rüzgar alanları kontrollü biçimde pasifleşir. Open-Meteo için ek anahtar gerekmez. Bu sayede proje hem demo hem üretim senaryosunda rahatça çalışabilir.",
    )
    add_bullet(doc, "WAQI / AQICN token: temel hava kalitesi akışı için gereklidir")
    add_bullet(doc, "Tomorrow.io key: gerçek zamanlı rüzgar özelliklerini güçlendirir")
    add_bullet(doc, "Open-Meteo: anahtarsız global fallback sağlar")
    add_bullet(doc, "Streamlit secrets yaklaşımı: güvenli dağıtım ve yerel geliştirme dengesi kurar")

    add_heading(doc, "11. Neden Güçlü Bir Proje?")
    add_bullet(doc, "Teknik olarak canlı veri, tahmin, analitik ve raporlamayı tek deneyimde birleştiriyor.")
    add_bullet(doc, "Ürün bakış açısından kullanıcıya sadece veri değil karar desteği sunuyor.")
    add_bullet(doc, "Küresel ölçekte genişlemeye uygun veri ve istasyon mimarisi kuruyor.")
    add_bullet(doc, "Profesyonel çıktılar sayesinde akademik, kurumsal ve sosyal paylaşım senaryolarına uyuyor.")
    add_bullet(doc, "Fallback stratejileri sayesinde servis eksikliklerinde bile kullanılabilir kalıyor.")

    add_heading(doc, "12. Sonuç")
    add_paragraph(
        doc,
        "AirPulse Global; çevresel veriyi sadeleştiren, anlamlandıran ve kullanıcıyı harekete geçiren bir platformdur. Proje, teknik olarak güçlü bir veri ve tahmin altyapısı üzerine kurulu olmasına rağmen, son kullanıcı açısından anlaşılır ve akıcı bir deneyim üretmeyi başarır. Bu yönüyle AirPulse Global, hem ürün düşüncesi hem de mühendislik disiplini açısından dengeli ve etkileyici bir projedir.",
    )

    footer_section = doc.sections[-1]
    footer = footer_section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer.add_run("AirPulse Global | Proje tanıtım dokümanı")
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = MUTED

    doc.save(DOCX_PATH)


def add_textbox(slide, left, top, width, height, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = PPTPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=20, color=PPT_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for item in items:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = PPTPt(font_size)
        run.font.color.rgb = color
        first = False
    return box


def add_card(slide, left, top, width, height, title, body, accent=PPT_BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPT_LIGHT
    shape.line.color.rgb = accent
    shape.line.width = PPTPt(1.5)
    add_textbox(slide, left + PPTInches(0.18), top + PPTInches(0.12), width - PPTInches(0.3), PPTInches(0.35), title, 18, accent, bold=True)
    add_textbox(slide, left + PPTInches(0.18), top + PPTInches(0.45), width - PPTInches(0.3), height - PPTInches(0.55), body, 12, PPT_DARK)


def add_header(slide, title, subtitle):
    add_textbox(slide, PPTInches(0.6), PPTInches(0.35), PPTInches(8.7), PPTInches(0.55), title, 28, PPT_DARK, bold=True)
    add_textbox(slide, PPTInches(0.6), PPTInches(0.9), PPTInches(10.6), PPTInches(0.35), subtitle, 12, PPT_MUTED)


def create_powerpoint(stats: dict[str, int | str]) -> None:
    prs = Presentation()
    prs.slide_width = PPTInches(13.333)
    prs.slide_height = PPTInches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PPT_LIGHT
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PPTInches(0.7),
        PPTInches(0.55),
        PPTInches(11.8),
        PPTInches(2.2),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = PPT_BLUE
    hero.line.color.rgb = PPT_PURPLE
    add_textbox(slide, PPTInches(1.05), PPTInches(1.05), PPTInches(5.6), PPTInches(0.6), "AirPulse Global", 26, PPTColor(255, 255, 255), bold=True)
    add_textbox(slide, PPTInches(1.05), PPTInches(1.6), PPTInches(6.6), PPTInches(0.35), "Professional Air Quality Intelligence", 15, PPTColor(230, 235, 255))
    add_textbox(slide, PPTInches(1.05), PPTInches(2.03), PPTInches(7.5), PPTInches(0.4), "Canlı hava verisi, tahmin, aksiyon önerisi ve raporlama tek platformda", 12, PPTColor(230, 235, 255))
    add_card(slide, PPTInches(0.95), PPTInches(3.25), PPTInches(2.45), PPTInches(1.15), "Canlı Veri", "AQI, PM2.5, PM10, NO2, O3 ve rüzgar bağlamı tek akışta")
    add_card(slide, PPTInches(3.62), PPTInches(3.25), PPTInches(2.45), PPTInches(1.15), "Tahmin", "Resmi WAQI forecast ve güvenli fallback yaklaşımı", accent=PPT_PURPLE)
    add_card(slide, PPTInches(6.29), PPTInches(3.25), PPTInches(2.45), PPTInches(1.15), "Aksiyon", "Kişisel sağlık ve sürdürülebilirlik önerileri", accent=PPT_GREEN)
    add_card(slide, PPTInches(8.96), PPTInches(3.25), PPTInches(2.45), PPTInches(1.15), "Raporlama", "PDF, JPEG ve CSV ile paylaşılabilir çıktı", accent=PPT_BLUE)
    add_textbox(slide, PPTInches(0.85), PPTInches(6.8), PPTInches(4.0), PPTInches(0.3), f"Hazırlanma tarihi: {stats['generated_at']}", 10, PPT_MUTED)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Projenin Kısa Tanımı", "AirPulse Global hangi problemi çözüyor?")
    add_textbox(
        slide,
        PPTInches(0.75),
        PPTInches(1.45),
        PPTInches(5.65),
        PPTInches(3.8),
        "AirPulse Global, açık hava kalitesi verisini yalnızca gösteren bir panel değil; bu veriyi yorumlayan, kullanıcı bağlamıyla birleştiren ve günlük kararları destekleyen bir çevresel zeka ürünüdür.",
        23,
        PPT_DARK,
    )
    add_bullets(
        slide,
        PPTInches(6.7),
        PPTInches(1.55),
        PPTInches(5.5),
        PPTInches(3.9),
        [
            "Dünya çapında izleme istasyonlarından veri toplar",
            "Rüzgar etkisini ve kirletici bağlamını gösterir",
            "Tahmin ve model açıklaması sunar",
            "Kullanıcıya uygulanabilir öneriler üretir",
            "Kurumsal sunuma uygun çıktılar hazırlar",
        ],
        font_size=19,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Showcase Snapshot", "Platformun dört temel değer alanı")
    snapshots = [
        ("Live Intelligence", "Gerçek zamanlı AQI, kirletici, istasyon ve rüzgar verisini tek akışta birleştirir."),
        ("Forecast Layer", "Observed-history yaklaşımı ve şeffaf model açıklamalarıyla yarın için görünürlük sağlar."),
        ("Action Engine", "Kişiselleştirilmiş günlük öneriler, puanlama ve hafif alışkanlık takibi üretir."),
        ("Export Ready", "A4 PDF, JPEG sosyal kart ve CSV çıktılarıyla sonucu paylaşılabilir hale getirir."),
    ]
    positions = [
        (PPTInches(0.8), PPTInches(1.55)),
        (PPTInches(6.9), PPTInches(1.55)),
        (PPTInches(0.8), PPTInches(4.0)),
        (PPTInches(6.9), PPTInches(4.0)),
    ]
    accents = [PPT_BLUE, PPT_PURPLE, PPT_GREEN, PPT_BLUE]
    for (title, body), (left, top), accent in zip(snapshots, positions, accents):
        add_card(slide, left, top, PPTInches(5.5), PPTInches(1.8), title, body, accent=accent)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Uygulama Nasıl Çalışıyor?", "Veri kaynağından kullanıcı çıktısına kadar akış")
    steps = [
        "1. Kullanıcı şehir veya istasyon seçer",
        "2. WAQI / Open-Meteo canlı veri akışı başlar",
        "3. Tomorrow.io varsa rüzgar bağlamı eklenir",
        "4. Cache katmanı performans ve kota uyumu sağlar",
        "5. Forecast sayfası resmi tahmini veya fallback modeli sunar",
        "6. Action Engine kişisel öneri üretir",
        "7. Reports katmanı sonuçları PDF, JPEG ve CSV olarak dışa aktarır",
    ]
    add_bullets(slide, PPTInches(0.9), PPTInches(1.5), PPTInches(6.0), PPTInches(4.6), steps, font_size=20)
    flow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(7.35), PPTInches(1.7), PPTInches(4.7), PPTInches(3.2))
    flow.fill.solid()
    flow.fill.fore_color.rgb = PPT_LIGHT
    flow.line.color.rgb = PPT_PURPLE
    add_textbox(slide, PPTInches(7.65), PPTInches(2.05), PPTInches(4.0), PPTInches(2.5), "Veri -> Analiz -> Tahmin -> Aksiyon -> Rapor\n\nAirPulse'un güçlü tarafı, bu beş halkayı tek ürün deneyiminde birbirine bağlamasıdır.", 20, PPT_DARK, bold=False, align=PP_ALIGN.CENTER)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Platform Özellikleri", "Kullanıcı deneyimini oluşturan ana modüller")
    add_card(slide, PPTInches(0.7), PPTInches(1.55), PPTInches(2.95), PPTInches(1.75), "Dashboard", "Şehir arama, canlı AQI kartları, bölgesel filtreler ve hızlı genel görünüm")
    add_card(slide, PPTInches(3.95), PPTInches(1.55), PPTInches(2.95), PPTInches(1.75), "Stations Map", "İstasyon yoğunluğu, kirletici popup'ları ve Forecast için kaynak seçimi", accent=PPT_PURPLE)
    add_card(slide, PPTInches(7.2), PPTInches(1.55), PPTInches(2.95), PPTInches(1.75), "Forecast", "Resmi günlük forecast ve güvenli geri dönüş modeli", accent=PPT_GREEN)
    add_card(slide, PPTInches(10.45), PPTInches(1.55), PPTInches(2.15), PPTInches(1.75), "Reports", "PDF, CSV, JPEG çıktılar", accent=PPT_BLUE)
    add_card(slide, PPTInches(1.8), PPTInches(4.05), PPTInches(4.0), PPTInches(1.8), "Take Action", "Karbon ayak izi, commute karşılaştırması, checklist ve streak takibi", accent=PPT_GREEN)
    add_card(slide, PPTInches(7.0), PPTInches(4.05), PPTInches(4.0), PPTInches(1.8), "Analytics", "Zaman desenleri, rüzgar dispersiyonu, karşılaştırmalar ve anomali okuması", accent=PPT_PURPLE)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Teknoloji Yığını", "Projenin arkasındaki ürün ve mühendislik bileşenleri")
    tech_points = [
        "Python 3.11+: modüler servisler ve temiz uygulama omurgası",
        "Streamlit: hızlı, çok sayfalı web arayüzü ve session state",
        "WAQI / AQICN: ana hava kalitesi kaynağı ve resmi forecast verisi",
        "Open-Meteo: anahtarsız global overview fallback",
        "Tomorrow.io: gerçek zamanlı rüzgar verisi",
        "Plotly: interaktif grafikler ve tahmin görselleştirmeleri",
        "Folium + Leaflet: etkileşimli istasyon haritası",
        "ReportLab + Pillow: PDF ve sosyal kart üretimi",
    ]
    add_bullets(slide, PPTInches(0.85), PPTInches(1.45), PPTInches(6.4), PPTInches(4.9), tech_points, font_size=19)
    metric_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(7.55), PPTInches(1.55), PPTInches(4.95), PPTInches(4.3))
    metric_box.fill.solid()
    metric_box.fill.fore_color.rgb = PPT_LIGHT
    metric_box.line.color.rgb = PPT_BLUE
    add_textbox(slide, PPTInches(7.9), PPTInches(1.9), PPTInches(4.2), PPTInches(0.35), "Repo Göstergeleri", 20, PPT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, PPTInches(8.0), PPTInches(2.55), PPTInches(4.0), PPTInches(2.8), f"{stats['modules']} Python modülü\n{stats['scripts']} yardımcı script\n{stats['notebooks']} araştırma notebook'u\n{stats['raw_city_files']} ham şehir veri dosyası\n{stats['stations']} istasyon kaydı", 22, PPT_DARK, align=PP_ALIGN.CENTER)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Forecast ve Güven Yaklaşımı", "Tahmin katmanında şeffaflık neden önemli?")
    add_textbox(slide, PPTInches(0.8), PPTInches(1.45), PPTInches(5.7), PPTInches(3.9), "Tahmin sisteminde amaç yalnızca bir sayı üretmek değil, kullanıcının gördüğü sonucun kaynağını dürüstçe anlatmaktır. Bu nedenle AirPulse, resmi sağlayıcı forecast ile kendi geri dönüş modelini net biçimde ayırır.", 22, PPT_DARK)
    add_bullets(slide, PPTInches(6.85), PPTInches(1.65), PPTInches(5.0), PPTInches(3.6), [
        "Öncelik: WAQI resmi günlük tahminleri",
        "Eksik veri durumunda: Holt-Winters fallback",
        "Arayüzde model kaynağı açıkça belirtilir",
        "Güven bantları belirsizliği görünür kılar",
        "Kullanıcı güveni korunur, veri sürekliliği bozulmaz",
    ], font_size=19)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Take Action ve Raporlama", "Veriyi kişisel karara ve paylaşılabilir çıktıya dönüştürmek")
    add_card(slide, PPTInches(0.85), PPTInches(1.55), PPTInches(5.6), PPTInches(2.05), "Take Action", "Sürüş, diyet, uçuş, elektrik ve ısıtma üzerinden karbon ayak izi hesaplar. Günlük checklist ve commute karşılaştırmalarıyla kullanıcıyı aksiyona taşır.", accent=PPT_GREEN)
    add_card(slide, PPTInches(6.85), PPTInches(1.55), PPTInches(5.6), PPTInches(2.05), "Reports", "A4 PDF, 1080x566 sosyal kart ve CSV çıktıları sayesinde teknik sonuçlar kurumsal veya sosyal paylaşım için hazır hale gelir.", accent=PPT_BLUE)
    add_textbox(slide, PPTInches(1.15), PPTInches(4.35), PPTInches(11.0), PPTInches(1.0), "Bu iki modül birlikte çalıştığında AirPulse yalnızca veri sunan bir uygulama olmaktan çıkar; kullanıcıya günlük karar desteği veren ve sonucu rapora dönüştüren tam bir ürün deneyimine dönüşür.", 22, PPT_DARK, align=PP_ALIGN.CENTER)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Mimari Güçlü Yanlar", "Projeyi teknik olarak öne çıkaran noktalar")
    add_bullets(slide, PPTInches(0.85), PPTInches(1.55), PPTInches(11.3), PPTInches(4.8), [
        "st.cache_data ile performans ve API kota uyumu birlikte korunuyor",
        "ThreadPoolExecutor ile çok şehirli veri akışı paralelleştiriliyor",
        "Session state ile kullanıcı bağlamı sayfalar arasında taşınıyor",
        "Modüler servis yapısı bakım kolaylığı ve yeniden kullanım sağlıyor",
        "Üçüncü taraf servis yoksa kontrollü fallback davranışı uygulanıyor",
        "Araştırma notebook'ları ile ürün kodu arasında net sorumluluk ayrımı bulunuyor",
    ], font_size=22)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Kullanım Senaryoları", "Bu proje nerelerde değer üretir?")
    add_card(slide, PPTInches(0.8), PPTInches(1.55), PPTInches(3.8), PPTInches(1.9), "Bireysel Kullanıcı", "Bugün dışarı çıkmalı mıyım, spor yapmalı mıyım, hangi önlemleri almalıyım?")
    add_card(slide, PPTInches(4.77), PPTInches(1.55), PPTInches(3.8), PPTInches(1.9), "Akademik / Demo", "Canlı veri, tahmin ve görselleştirme yetkinliklerini tek projede göstermek için")
    add_card(slide, PPTInches(8.74), PPTInches(1.55), PPTInches(3.8), PPTInches(1.9), "Kurumsal İletişim", "Paylaşılabilir PDF ve sosyal kartlarla çevre verisini anlaşılır sunmak", accent=PPT_PURPLE)
    add_card(slide, PPTInches(2.8), PPTInches(4.1), PPTInches(7.6), PPTInches(1.75), "Karar Destek Perspektifi", "AirPulse'un asıl gücü, veriyi teknik bir panel olarak bırakmaması; kullanıcı davranışı, sürdürülebilirlik ve iletişim katmanlarıyla birleştirmesidir.", accent=PPT_GREEN)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Sonuç", "AirPulse Global neden etkileyici bir proje?")
    add_textbox(slide, PPTInches(0.9), PPTInches(1.65), PPTInches(11.6), PPTInches(1.2), "AirPulse Global, hava kalitesi verisini gerçek zamanlı izleyen, tahmin eden, açıklayan, kullanıcıyı aksiyona yönlendiren ve sonucu profesyonel çıktılara dönüştüren bütünleşik bir platformdur.", 24, PPT_DARK, align=PP_ALIGN.CENTER)
    add_bullets(slide, PPTInches(2.0), PPTInches(3.3), PPTInches(9.0), PPTInches(2.2), [
        "Teknik olarak güçlü",
        "Ürün dili açısından anlaşılır",
        "Sunum ve demo için etkili",
        "Gelecek geliştirmelere açık",
    ], font_size=24)
    add_textbox(slide, PPTInches(3.1), PPTInches(6.5), PPTInches(7.2), PPTInches(0.35), "AirPulse Global | Environmental intelligence made understandable", 14, PPT_MUTED, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


def main() -> None:
    stats = project_stats()
    create_word_document(stats)
    create_powerpoint(stats)
    print(f"Created: {DOCX_PATH.name}")
    print(f"Created: {PPTX_PATH.name}")


if __name__ == "__main__":
    main()
